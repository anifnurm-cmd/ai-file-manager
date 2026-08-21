from __future__

import hashlib
import json
import os
import re
import sqlite3
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.error import URLError
from urllib.request import Request, urlopen

from fastapi import FastAPI, HTTPException
from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field

APP_ROOT = Path(__file__).resolve().parent
DATA_ROOT = Path(os.getenv("AFM_DATA_DIR", str(APP_ROOT / "data"))).resolve()
DB_PATH = DATA_ROOT / "index.sqlite3"
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434").rstrip("/")
OLLAMA_CHAT_MODEL = os.getenv("OLLAMA_CHAT_MODEL", "qwen3:8b")
OLLAMA_EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")
SUPPORTED = {
    ".txt", ".md", ".csv", ".json", ".xml", ".log",
    ".pdf", ".docx", ".xlsx", ".xlsm", ".pptx",
    ".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff",
}

DATA_ROOT.mkdir(parents=True, exist_ok=True)


def db() -> sqlite3.Connection:
    c = sqlite3.connect(DB_PATH)
    c.row_factory = sqlite3.Row
    return c


def init_db() -> None:
    with db() as c:
        c.executescript(
            """
            PRAGMA journal_mode=WAL;
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY,
                path TEXT UNIQUE NOT NULL,
                name TEXT NOT NULL,
                ext TEXT NOT NULL,
                size INTEGER NOT NULL,
                mtime REAL NOT NULL,
                sha256 TEXT NOT NULL,
                title TEXT NOT NULL,
                doc_type TEXT NOT NULL,
                summary TEXT NOT NULL,
                content TEXT NOT NULL,
                indexed_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS chunks (
                id INTEGER PRIMARY KEY,
                document_id INTEGER NOT NULL REFERENCES documents(id) ON DELETE CASCADE,
                chunk_no INTEGER NOT NULL,
                text TEXT NOT NULL,
                embedding TEXT,
                UNIQUE(document_id, chunk_no)
            );
            CREATE VIRTUAL TABLE IF NOT EXISTS documents_fts USING fts5(
                document_id UNINDEXED,
                path,
                name,
                title,
                doc_type,
                summary,
                content
            );
            CREATE INDEX IF NOT EXISTS idx_documents_sha ON documents(sha256);
            CREATE INDEX IF NOT EXISTS idx_chunks_document ON chunks(document_id);
            """
        )


init_db()


@dataclass
class ParsedDocument:
    content: str
    title: str
    doc_type: str
    summary: str


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def split_chunks(text: str, size: int = 1200, overlap: int = 150) -> list[str]:
    clean = text.strip()
    if not clean:
        return []
    if len(clean) <= size:
        return [clean]
    out: list[str] = []
    start = 0
    while start < len(clean):
        end = min(start + size, len(clean))
        piece = clean[start:end].strip()
        if piece:
            out.append(piece)
        if end >= len(clean):
            break
        start = max(end - overlap, start + 1)
    return out


def classify(path: Path, text: str) -> tuple[str, str]:
    low = f"{path.name} {text[:6000]}".lower()
    rules = [
        ("surat_undangan", ("undangan", "surat undangan")),
        ("surat", ("surat", "nomor surat", "perihal")),
        ("proposal", ("proposal", "latar belakang", "rencana kegiatan")),
        ("laporan", ("laporan", "hasil kegiatan", "kesimpulan")),
        ("notulen", ("notulen", "notula", "hasil rapat")),
        ("data_peserta_didik", ("nisn", "peserta didik", "nama siswa")),
        ("jadwal", ("jadwal", "tanggal", "pukul", "tempat")),
    ]
    for name, needles in rules:
        if any(n in low for n in needles):
            return name, make_title(path, text)
    return "dokumen", make_title(path, text)


def make_title(path: Path, text: str) -> str:
    lines = [normalize(x) for x in text.splitlines() if normalize(x)]
    for line in lines[:40]:
        if 8 <= len(line) <= 160 and not re.fullmatch(r"[\W_\d]+", line):
            return line
    return re.sub(r"[_-]+", " ", path.stem).strip() or path.name


def summarize(text: str) -> str:
    clean = normalize(text)
    if not clean:
        return "No extractable text."
    sentences = re.split(r"(?<=[.!?])\s+", clean)
    return " ".join(sentences[:4])[:900]


def parse_file(path: Path) -> ParsedDocument:
    ext = path.suffix.lower()
    text = ""
    if ext in {".txt", ".md", ".csv", ".json", ".xml", ".log"}:
        text = path.read_text(encoding="utf-8", errors="ignore")
    elif ext == ".pdf":
        import fitz
        with fitz.open(path) as pdf:
            text = "\n".join(page.get_text("text") for page in pdf)
    elif ext == ".docx":
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = "\n".join(parts)
    elif ext in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"[SHEET] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if any(values):
                    parts.append(" | ".join(values))
        text = "\n".join(parts)
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                parts.append(f"[SLIDE {idx}] " + " | ".join(slide_text))
        text = "\n".join(parts)
    elif ext in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tif", ".tiff"}:
        try:
            import pytesseract
            from PIL import Image
            text = pytesseract.image_to_string(Image.open(path))
        except Exception:
            text = ""
    else:
        raise ValueError(f"Unsupported extension: {ext}")

    content = text.strip()
    doc_type, title = classify(path, content)
    return ParsedDocument(content=content, title=title, doc_type=doc_type, summary=summarize(content))


def ollama_json(endpoint: str, payload: dict[str, Any], timeout: int = 90) -> dict[str, Any] | None:
    try:
        req = Request(
            f"{OLLAMA_URL}{endpoint}",
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urlopen(req, timeout=timeout) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except (OSError, URLError, TimeoutError, json.JSONDecodeError):
        return None


def ollama_available() -> bool:
    try:
        with urlopen(Request(f"{OLLAMA_URL}/api/tags"), timeout=2):
            return True
    except Exception:
        return False


def embed(text: str) -> list[float] | None:
    if not text.strip():
        return None
    data = ollama_json("/api/embed", {"model": OLLAMA_EMBED_MODEL, "input": text[:6000]}, timeout=30)
    if data:
        embeddings = data.get("embeddings")
        if isinstance(embeddings, list) and embeddings:
            first = embeddings[0]
            if isinstance(first, list) and all(isinstance(x, (int, float)) for x in first):
                return [float(x) for x in first]
    data = ollama_json("/api/embeddings", {"model": OLLAMA_EMBED_MODEL, "prompt": text[:6000]}, timeout=30)
    if data and isinstance(data.get("embedding"), list):
        return [float(x) for x in data["embedding"]]
    return None


def cosine(a: list[float], b: list[float]) -> float:
    import math
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x * y for x, y in zip(a, b))
    na = math.sqrt(sum(x * x for x in a))
    nb = math.sqrt(sum(y * y for y in b))
    return dot / (na * nb) if na and nb else 0.0


def extract_ai_metadata(parsed: ParsedDocument) -> ParsedDocument:
    if not ollama_available() or not parsed.content:
        return parsed
    prompt = (
        "Return JSON only with keys title, doc_type, summary. "
        "Do not invent facts. Keep title <= 160 chars and summary <= 900 chars. "
        f"Current title: {parsed.title}\nCurrent type: {parsed.doc_type}\n"
        f"CONTENT:\n{parsed.content[:8000]}"
    )
    response = ollama_json("/api/generate", {"model": OLLAMA_CHAT_MODEL, "prompt": prompt, "stream": False}, timeout=90)
    if not response or not isinstance(response.get("response"), str):
        return parsed
    match = re.search(r"\{.*\}", response["response"], re.S)
    if not match:
        return parsed
    try:
        data = json.loads(match.group(0))
        return ParsedDocument(
            content=parsed.content,
            title=str(data.get("title") or parsed.title)[:160],
            doc_type=str(data.get("doc_type") or parsed.doc_type)[:80],
            summary=str(data.get("summary") or parsed.summary)[:900],
        )
    except (json.JSONDecodeError, TypeError):
        return parsed


def upsert_document(path: Path) -> tuple[int, str]:
    stat = path.stat()
    digest = sha256_file(path)
    with db() as c:
        existing = c.execute("SELECT id, sha256 FROM documents WHERE path = ?", (str(path),)).fetchone()
        if existing and existing["sha256"] == digest:
            return int(existing["id"]), "skipped"

    parsed = extract_ai_metadata(parse_file(path))
    chunks = split_chunks(parsed.content)
    embeddings = [embed(x) for x in chunks] if ollama_available() else [None] * len(chunks)

    with db() as c:
        c.execute("PRAGMA foreign_keys=ON")
        old = c.execute("SELECT id FROM documents WHERE path = ?", (str(path),)).fetchone()
        if old:
            doc_id = int(old["id"])
            c.execute("DELETE FROM documents_fts WHERE document_id = ?", (doc_id,))
            c.execute("DELETE FROM chunks WHERE document_id = ?", (doc_id,))
            c.execute(
                """UPDATE documents SET name=?, ext=?, size=?, mtime=?, sha256=?, title=?, doc_type=?, summary=?, content=?, indexed_at=CURRENT_TIMESTAMP WHERE id=?""",
                (path.name, path.suffix.lower(), stat.st_size, stat.st_mtime, digest, parsed.title, parsed.doc_type, parsed.summary, parsed.content, doc_id),
            )
        else:
            cur = c.execute(
                """INSERT INTO documents(path,name,ext,size,mtime,sha256,title,doc_type,summary,content) VALUES(?,?,?,?,?,?,?,?,?,?)""",
                (str(path), path.name, path.suffix.lower(), stat.st_size, stat.st_mtime, digest, parsed.title, parsed.doc_type, parsed.summary, parsed.content),
            )
            doc_id = int(cur.lastrowid)
        c.execute(
            "INSERT INTO documents_fts(document_id,path,name,title,doc_type,summary,content) VALUES(?,?,?,?,?,?,?)",
            (doc_id, str(path), path.name, parsed.title, parsed.doc_type, parsed.summary, parsed.content),
        )
        for i, chunk in enumerate(chunks):
            c.execute("INSERT INTO chunks(document_id,chunk_no,text,embedding) VALUES(?,?,?,?)", (doc_id, i, chunk, json.dumps(embeddings[i]) if embeddings[i] is not None else None))
    return doc_id, "indexed"


def keyword_results(query: str, limit: int = 20) -> list[dict[str, Any]]:
    terms = [x for x in re.findall(r"[\w-]+", query, re.UNICODE) if len(x) > 1][:10]
    if not terms:
        return []
    expression = " OR ".join('"' + t.replace('"', "") + '"*' for t in terms)
    try:
        with db() as c:
            return [dict(r) for r in c.execute(
                """SELECT d.id,d.name,d.path,d.title,d.doc_type,d.summary,bm25(documents_fts) AS score
                FROM documents_fts f JOIN documents d ON d.id=f.document_id
                WHERE documents_fts MATCH ? ORDER BY score LIMIT ?""",
                (expression, max(1, min(limit, 50))),
            )]
    except sqlite3.OperationalError:
        return []


def semantic_results(query: str, limit: int = 8) -> list[dict[str, Any]]:
    qv = embed(query)
    if not qv:
        return []
    scored: dict[int, tuple[float, dict[str, Any]]] = {}
    with db() as c:
        rows = c.execute(
            """SELECT ch.id,ch.document_id,ch.text,ch.embedding,d.name,d.path,d.title,d.doc_type,d.summary
            FROM chunks ch JOIN documents d ON d.id=ch.document_id WHERE ch.embedding IS NOT NULL"""
        ).fetchall()
    for r in rows:
        try:
            ev = json.loads(r["embedding"])
            score = cosine(qv, ev)
        except Exception:
            continue
        item = {
            "id": r["document_id"], "name": r["name"], "path": r["path"], "title": r["title"],
            "doc_type": r["doc_type"], "summary": r["summary"], "chunk": r["text"], "score": score,
        }
        previous = scored.get(item["id"])
        if previous is None or score > previous[0]:
            scored[item["id"]] = (score, item)
    return [x[1] for x in sorted(scored.values(), key=lambda x: x[0], reverse=True)[:limit]]


def hybrid_results(query: str, limit: int = 8) -> list[dict[str, Any]]:
    sem = semantic_results(query, limit * 2)
    kw = keyword_results(query, limit * 2)
    merged: dict[int, dict[str, Any]] = {}
    for rank, item in enumerate(sem):
        merged[item["id"]] = {**item, "hybrid": 1.0 / (rank + 1)}
    for rank, item in enumerate(kw):
        if item["id"] in merged:
            merged[item["id"]]["hybrid"] += 0.8 / (rank + 1)
        else:
            merged[item["id"]] = {**item, "hybrid": 0.8 / (rank + 1)}
    return sorted(merged.values(), key=lambda x: x["hybrid"], reverse=True)[:limit]


def generate_answer(question: str, hits: list[dict[str, Any]]) -> str:
    if not hits:
        return "I could not find relevant indexed documents."
    context = "\n\n".join(
        f"FILE: {h['name']}\nPATH: {h['path']}\nTITLE: {h['title']}\nTYPE: {h['doc_type']}\n"
        + (f"MATCHED PASSAGE: {h.get('chunk','')}\n" if h.get("chunk") else f"SUMMARY: {h['summary']}\n")
        for h in hits
    )
    if ollama_available():
        prompt = (
            "Answer ONLY from the supplied file context. Do not invent facts. "
            "When uncertain, say so. Mention relevant filenames. Reply in the same language as the question.\n\n"
            f"QUESTION:\n{question}\n\nCONTEXT:\n{context}"
        )
        response = ollama_json("/api/generate", {"model": OLLAMA_CHAT_MODEL, "prompt": prompt, "stream": False}, timeout=120)
        if response and isinstance(response.get("response"), str) and response["response"].strip():
            return response["response"].strip()
    return "Relevant files found:\n" + "\n".join(f"- {h['title']} — {h['path']}" for h in hits)


class ScanRequest(BaseModel):
    path: str = Field(min_length=1)

class SearchRequest(BaseModel):
    query: str = Field(min_length=1)
    limit: int = Field(default=20, ge=1, le=50)

class AskRequest(BaseModel):
    question: str = Field(min_length=1)

class RenameRequest(BaseModel):
    id: int
    new_name: str = Field(min_length=1, max_length=240)


app = FastAPI(title="AI File Manager", version="1.0.0")

@app.get("/", response_class=HTMLResponse)
def home() -> str:
    return HTML

@app.get("/health")
def health() -> dict[str, Any]:
    return {"ok": True, "version": app.version, "ollama": ollama_available()}

@app.get("/api/status")
def status() -> dict[str, Any]:
    with db() as c:
        count = c.execute("SELECT COUNT(*) FROM documents").fetchone()[0]
    return {"documents": int(count), "ollama": ollama_available(), "chat_model": OLLAMA_CHAT_MODEL, "embed_model": OLLAMA_EMBED_MODEL}

@app.post("/api/scan")
def scan(req: ScanRequest) -> dict[str, Any]:
    root = Path(req.path).expanduser().resolve()
    if not root.is_dir():
        raise HTTPException(400, "Folder does not exist or is not accessible")
    seen = indexed = skipped = errors = 0
    error_details: list[str] = []
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED:
            continue
        seen += 1
        try:
            _, state = upsert_document(path)
            if state == "indexed": indexed += 1
            else: skipped += 1
        except Exception as exc:
            errors += 1
            if len(error_details) < 20: error_details.append(f"{path}: {exc}")
    return {"seen": seen, "indexed": indexed, "skipped": skipped, "errors": errors, "error_details": error_details}

@app.post("/api/search")
def search(req: SearchRequest) -> dict[str, Any]:
    return {"results": hybrid_results(req.query, req.limit)}

@app.post("/api/ask")
def ask(req: AskRequest) -> dict[str, Any]:
    hits = hybrid_results(req.question, 8)
    answer = generate_answer(req.question, hits)
    sources = [{k: h.get(k) for k in ("id", "name", "path", "title", "doc_type", "summary")} for h in hits]
    return {"answer": answer, "sources": sources, "semantic": bool(semantic_results(req.question, 1))}

@app.post("/api/rename")
def rename(req: RenameRequest) -> dict[str, Any]:
    safe_name = Path(req.new_name).name
    if safe_name != req.new_name or any(ch in safe_name for ch in '<>:"/\\|?*'):
        raise HTTPException(400, "Invalid Windows filename")
    with db() as c:
        row = c.execute("SELECT path FROM documents WHERE id = ?", (req.id,)).fetchone()
    if not row: raise HTTPException(404, "Document not found in index")
    old = Path(row["path"])
    if not old.exists(): raise HTTPException(404, "Original file no longer exists")
    new = old.with_name(safe_name)
    if new.exists() and new.resolve() != old.resolve(): raise HTTPException(409, "Target filename already exists")
    old.rename(new)
    upsert_document(new)
    with db() as c:
        c.execute("DELETE FROM documents WHERE path = ?", (str(old),))
        c.execute("DELETE FROM documents_fts WHERE path = ?", (str(old),))
    return {"ok": True, "new_path": str(new)}


HTML = r'''<!doctype html>
<html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>AI File Manager</title>
<style>
:root{font-family:Inter,system-ui,sans-serif;color:#e8edf5;background:#09111d}*{box-sizing:border-box}body{margin:0;background:radial-gradient(circle at top,#182b46,#09111d 58%);min-height:100vh}.wrap{max-width:1100px;margin:auto;padding:28px 18px}.top{display:flex;justify-content:space-between;align-items:flex-start;gap:16px}.eyebrow{color:#8ea2bd;letter-spacing:.12em;font-size:12px}.title{font-size:42px;margin:5px 0 8px}.muted{color:#91a4bc}.status{border:1px solid #314862;border-radius:999px;padding:8px 12px}.card{background:#101b2aee;border:1px solid #253a53;border-radius:18px;padding:18px;margin-top:16px}.row{display:flex;gap:10px}.grow{flex:1}input,textarea{width:100%;background:#0a1421;color:#edf4ff;border:1px solid #304861;border-radius:10px;padding:12px;font:inherit}button{background:#77b9ff;color:#08121d;border:0;border-radius:10px;padding:12px 16px;font-weight:800;cursor:pointer}.grid{display:grid;gap:10px;margin-top:12px}.item{display:flex;justify-content:space-between;gap:12px;padding:13px;border:1px solid #263d55;background:#0b1523;border-radius:12px}.item b{display:block;margin-bottom:4px}.small{font-size:12px;color:#7f95ae}.answer{white-space:pre-wrap;line-height:1.65;margin-top:12px}.err{color:#ffaaa7}.source{margin-top:8px;color:#a8b9cd}@media(max-width:700px){.top,.row{flex-direction:column}.title{font-size:32px}}
</style></head><body><main class="wrap"><section class="top"><div><div class="eyebrow">LOCAL-FIRST DOCUMENT AI</div><div class="title">AI File Manager</div><div class="muted">Index mixed documents, recall their contents, ask questions, and rename safely.</div></div><div id="status" class="status">Checking…</div></section>
<section class="card"><h2>Index a folder</h2><div class="row"><input class="grow" id="folder" placeholder="C:\\Users\\You\\Documents"><button onclick="scan()">Scan</button></div><div id="scanmsg" class="muted" style="margin-top:10px"></div></section>
<section class="card"><h2>Ask your files</h2><div class="row"><textarea class="grow" id="question" rows="3" placeholder="Which documents discuss Renstra? Summarize the proposal files."></textarea><button onclick="ask()">Ask</button></div><div id="answer" class="answer muted">No question yet.</div><div id="sources"></div></section>
<section class="card"><h2>Search files</h2><div class="row"><input class="grow" id="query" placeholder="Search by content, title, summary, type"><button onclick="searchFiles()">Search</button></div><div id="results" class="grid"></div></section>
<p class="muted">Core indexing and search are local. Ollama adds semantic retrieval, AI metadata, and natural-language answers. OCR is optional.</p></main>
<script>
const $=id=>document.getElementById(id);async function api(url,opts={}){const r=await fetch(url,{headers:{'Content-Type':'application/json'},...opts});const d=await r.json();if(!r.ok)throw Error(d.detail||'Request failed');return d}
async function refresh(){const d=await api('/api/status');$('status').textContent=d.ollama?'AI online · '+d.documents+' docs':'Local mode · '+d.documents+' docs'}
async function scan(){try{$('scanmsg').textContent='Scanning and indexing…';const d=await api('/api/scan',{method:'POST',body:JSON.stringify({path:$('folder').value})});$('scanmsg').textContent=`Seen ${d.seen} · indexed ${d.indexed} · skipped ${d.skipped} · errors ${d.errors}`;if(d.error_details.length)$('scanmsg').innerHTML+=`<div class="err">${esc(d.error_details.join('\n'))}</div>`;refresh()}catch(e){$('scanmsg').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}
async function searchFiles(){const q=$('query').value.trim();if(!q)return;try{const d=await api('/api/search',{method:'POST',body:JSON.stringify({query:q,limit:20})});$('results').innerHTML=d.results.map(x=>`<div class="item"><div><b>${esc(x.title||x.name)}</b><div class="small">${esc(x.path)} · ${esc(x.doc_type)}</div><div class="muted">${esc(x.summary||'')}</div></div><button onclick='renameFile(${x.id},${JSON.stringify(x.name)})'>Rename</button></div>`).join('')||'<div class="muted">No results.</div>'}catch(e){$('results').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}
async function ask(){const q=$('question').value.trim();if(!q)return;$('answer').textContent='Searching your files…';$('sources').innerHTML='';try{const d=await api('/api/ask',{method:'POST',body:JSON.stringify({question:q})});$('answer').textContent=d.answer;$('sources').innerHTML=d.sources.map(x=>`<div class="source">• ${esc(x.title)} — ${esc(x.path)}</div>`).join('')}catch(e){$('answer').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}
async function renameFile(id,old){const n=prompt('New filename',old);if(!n)return;try{await api('/api/rename',{method:'POST',body:JSON.stringify({id,new_name:n})});await searchFiles();await refresh()}catch(e){alert(e.message)}}
function esc(s){return String(s??'').replace(/[&<>"']/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]))}refresh();
</script></body></html>'''

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="127.0.0.1", port=int(os.getenv("PORT", "8787")))
