from __future__ import annotations

import hashlib
import re
from pathlib import Path

from fastapi import FastAPI, HTTPException
from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field

ROOT = Path(__file__).resolve().parent
SUPPORTED = {'.pdf','.docx','.xlsx','.xlsm','.pptx','.txt','.md','.csv','.json','.xml','.log','.jpg','.jpeg','.png','.webp','.bmp','.tif','.tiff'}


def sha(path: Path) -> str:
    h = hashlib.sha256()
    with path.open('rb') as f:
        for block in iter(lambda: f.read(1024 * 1024), b''):
            h.update(block)
    return h.hexdigest()


def text(path: Path) -> str:
    e = path.suffix.lower()
    if e in {'.txt','.md','.csv','.json','.xml','.log'}:
        return path.read_text(encoding='utf-8', errors='ignore')
    if e == '.pdf':
        import fitz
        with fitz.open(path) as pdf:
            return '\n'.join(page.get_text('text') for page in pdf)
    if e == '.docx':
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(' | '.join(cell.text.strip() for cell in row.cells))
        return '\n'.join(parts)
    if e in {'.xlsx','.xlsm'}:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts=[]
        for ws in wb.worksheets:
            parts.append(f'[SHEET] {ws.title}')
            for row in ws.iter_rows(values_only=True):
                vals=['' if v is None else str(v) for v in row]
                if any(vals): parts.append(' | '.join(vals))
        return '\n'.join(parts)
    if e == '.pptx':
        from pptx import Presentation
        prs=Presentation(path); parts=[]
        for i,slide in enumerate(prs.slides,1):
            vals=[shape.text.strip() for shape in slide.shapes if hasattr(shape,'text') and shape.text.strip()]
            if vals: parts.append(f'[SLIDE {i}] ' + ' | '.join(vals))
        return '\n'.join(parts)
    if e in {'.jpg','.jpeg','.png','.webp','.bmp','.tif','.tiff'}:
        try:
            import pytesseract
            from PIL import Image
            return pytesseract.image_to_string(Image.open(path))
        except Exception:
            return ''
    raise ValueError(f'Unsupported file type: {e}')


def clean(s: str) -> str:
    return re.sub(r'\s+', ' ', s or '').strip()


def metadata(path: Path, body: str):
    lines=[clean(x) for x in body.splitlines() if clean(x)]
    title=next((x for x in lines[:50] if 8<=len(x)<=160 and not re.fullmatch(r'[\W_\d]+',x)), path.stem.replace('_',' ').replace('-',' '))
    low=f'{path.name} {body[:7000]}'.lower()
    rules=[
      ('surat_undangan',['surat undangan','undangan']),
      ('proposal',['proposal','latar belakang','rencana kegiatan']),
      ('notulen',['notulen','notula','hasil rapat']),
      ('laporan',['laporan','hasil kegiatan','kesimpulan']),
      ('data_peserta_didik',['nisn','peserta didik','nama siswa']),
      ('jadwal',['jadwal','pukul','hari/tanggal']),
      ('surat',['nomor surat','perihal','sehubungan dengan'])]
    doc_type='dokumen'
    for kind, needles in rules:
        if any(x in low for x in needles):
            doc_type=kind
            break
    sentences=re.split(r'(?<=[.!?])\s+',clean(body))
    summary=' '.join(sentences[:4])[:900] or 'Tidak ada teks yang dapat diekstrak.'
    return title[:160], doc_type, summary


def chunk_text(body: str, size=900, overlap=120):
    body=body.strip()
    if not body: return []
    if len(body)<=size: return [body]
    out=[]
    start=0
    while start<len(body):
        end=min(start+size,len(body))
        piece=body[start:end].strip()
        if piece: out.append(piece)
        if end>=len(body): break
        start=max(start+1,end-overlap)
    return out


class ScanRequest(BaseModel):
    path: str = Field(min_length=1)
class Query(BaseModel):
    query: str = Field(min_length=1)
    limit: int = Field(default=20, ge=1, le=50)
class Ask(BaseModel):
    question: str = Field(min_length=1)
class Rename(BaseModel):
    id: int
    new_name: str = Field(min_length=1, max_length=240)


app=FastAPI(title='AI File Manager',version='3.0.0')

def index(path: Path):
    raise RuntimeError('Semantic runner is not loaded')

def search(q: Query):
    return {'results': [], 'semantic_ready': False}

def ask(q: Ask):
    return {'answer':'Semantic runner is not loaded.','sources':[]}

HTML='''<!doctype html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>AI File Manager</title><style>body{margin:0;background:#08121e;color:#eaf2f8;font:15px system-ui}.wrap{max-width:1100px;margin:auto;padding:28px 18px}.top{display:flex;justify-content:space-between}.card{background:#101d2b;border:1px solid #29435d;border-radius:18px;padding:18px;margin-top:16px}.row{display:flex;gap:10px}.grow{flex:1}input,textarea{width:100%;background:#091521;color:#edf4fb;border:1px solid #345069;border-radius:10px;padding:12px;font:inherit}button{background:#76b9f8;color:#06111b;border:0;border-radius:10px;padding:12px 16px;font-weight:800}.item{display:flex;justify-content:space-between;gap:12px;padding:12px;border:1px solid #29415b;border-radius:12px;margin-top:10px}.muted{color:#90a7be}.err{color:#ffaaa7}.source{margin-top:7px;color:#9db4ca}.answer{white-space:pre-wrap;line-height:1.6;margin-top:12px}@media(max-width:700px){.top,.row,.item{flex-direction:column}}</style></head><body><main class="wrap"><section class="top"><div><div>LOCAL-FIRST DOCUMENT AI</div><h1>AI File Manager</h1><div class="muted">Semantic recall + AI reranking + safe rename</div></div><div id="status">Checking…</div></section><section class="card"><h2>Index a folder</h2><div class="row"><input id="folder" class="grow" placeholder="C:\\Users\\You\\Documents"><button onclick="scan()">Scan</button><button onclick="reindex()">Rebuild semantic index</button></div><div id="scanmsg" class="muted"></div></section><section class="card"><h2>Ask your files</h2><div class="row"><textarea id="q" class="grow" rows="3" placeholder="Dokumen apa yang membahas Renstra? Ringkas proposal sarpras."></textarea><button onclick="ask()">Ask</button></div><div id="answer" class="answer muted">Belum ada pertanyaan.</div><div id="sources"></div></section><section class="card"><h2>Search</h2><div class="row"><input id="s" class="grow" placeholder="Cari berdasarkan makna atau isi dokumen"><button onclick="searchFiles()">Search</button></div><div id="results"></div></section></main><script>const $=x=>document.getElementById(x);async function api(u,o={}){const r=await fetch(u,{headers:{'Content-Type':'application/json'},...o});const d=await r.json();if(!r.ok)throw Error(d.detail||'Request failed');return d}async function status(){const d=await api('/api/status');$('status').textContent=d.semantic_ready?`Semantic ready · ${d.embedded_documents}/${d.documents}`:`Semantic unavailable · ${d.documents} docs`}function esc(s){return String(s??'').replace(/[&<>"']/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]))}function render(x){return `<div class="item"><div><b>${esc(x.title||x.name)}</b><div class="muted">${esc(x.path)}</div><div>${esc(x.summary||'')} · ${Number(x.relevance||0).toFixed(0)}% relevant</div></div><button onclick='renameFile(${x.id},${JSON.stringify(x.name)})'>Rename</button></div>`}async function scan(){try{$('scanmsg').textContent='Scanning, extracting and embedding…';const d=await api('/api/scan',{method:'POST',body:JSON.stringify({path:$('folder').value})});$('scanmsg').textContent=`Seen ${d.seen} · indexed ${d.indexed} · skipped ${d.skipped} · errors ${d.errors}`;if(d.error_details.length)$('scanmsg').innerHTML+='<div class="err">'+esc(d.error_details.join('\\n'))+'</div>';status()}catch(e){$('scanmsg').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}async function reindex(){try{$('scanmsg').textContent='Rebuilding semantic index…';const d=await api('/api/reindex-semantic',{method:'POST'});$('scanmsg').textContent=`Reindexed ${d.reindexed}/${d.documents} · errors ${d.errors}`;status()}catch(e){$('scanmsg').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}async function searchFiles(){const q=$('s').value.trim();if(!q)return;try{const d=await api('/api/search',{method:'POST',body:JSON.stringify({query:q,limit:20})});$('results').innerHTML=d.results.map(render).join('')||'<div class="muted">No sufficiently relevant results.</div>'}catch(e){$('results').innerHTML='<span class="err">'+esc(e.message)+'</span>'}}async function ask(){const q=$('q').value.trim();if(!q)return;$('answer').textContent='Finding relevant passages and asking AI…';try{const d=await api('/api/ask',{method:'POST',body:JSON.stringify({question:q})});$('answer').textContent=d.answer;$('sources').innerHTML=d.sources.map(x=>`<div class="source">• ${esc(x.title)} — ${esc(x.path)} · ${Number(x.relevance||0).toFixed(0)}%</div>`).join('')}catch(e){$('answer').textContent=e.message}}async function renameFile(id,old){const n=prompt('New filename',old);if(!n)return;try{await api('/api/rename',{method:'POST',body:JSON.stringify({id,new_name:n})});await searchFiles();await status()}catch(e){alert(e.message)}}status();</script></body></html>'''

@app.get('/',response_class=HTMLResponse)
def home(): return HTML

@app.get('/health')
def health(): return {'ok':True,'version':app.version}

@app.get('/api/status')
def status(): return {'documents':0,'embedded_documents':0,'semantic_ready':False}

@app.post('/api/scan')
def scan(req:ScanRequest):
    root=Path(req.path).expanduser().resolve()
    if not root.is_dir(): raise HTTPException(400,'Folder does not exist or is not accessible')
    seen=indexed=skipped=errors=0
    details=[]
    for p in root.rglob('*'):
        if not p.is_file() or p.suffix.lower() not in SUPPORTED: continue
        seen+=1
        try:
            _,state=index(p); indexed+=state=='indexed'; skipped+=state=='skipped'
        except Exception as e:
            errors+=1
            if len(details)<20: details.append(f'{p}: {e}')
    return {'seen':seen,'indexed':indexed,'skipped':skipped,'errors':errors,'error_details':details}

@app.post('/api/search')
def do_search(req:Query): return search(req)

@app.post('/api/ask')
def do_ask(req:Ask): return ask(req)

@app.post('/api/reindex-semantic')
def reindex(): return {'documents':0,'reindexed':0,'errors':0}

@app.post('/api/rename')
def rename(req:Rename):
    raise HTTPException(501,'Rename is provided by the semantic runner')
